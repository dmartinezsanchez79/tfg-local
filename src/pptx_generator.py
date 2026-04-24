"""Generación de presentaciones PPTX desde la KnowledgeBase.

Pipeline (por átomos):
1. `plan_slides`    : el LLM asigna átomos de la KB a slides tipadas.
2. `_render_bullets`: una llamada por slide, prompt por `kind`.
3. `_render_bullets`: también genera la conclusión final.
4. `render_pptx`    : montaje determinista sobre la plantilla institucional.
"""
from __future__ import annotations

import logging
import re
import unicodedata
from dataclasses import dataclass, field
from io import BytesIO
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Pt
from pydantic import BaseModel, Field, ValidationError

from .config import (
    DEFAULT_NUM_SLIDES_MAX,
    DEFAULT_NUM_SLIDES_MIN,
    LAYOUT_CONTENT,
    LAYOUT_TITLE,
    MAX_BULLETS_PER_SLIDE,
    MAX_CHARS_PER_BULLET,
    MAX_CHARS_SLIDE_TITLE,
    TEMPLATE_PATH,
)
from .exceptions import GenerationError, TemplateError
from .knowledge_base import (
    Definition,
    Example,
    FormulaOrCode,
    KnowledgeBase,
    NumericDatum,
    Relation,
    relation_to_natural,
)
from .map_reduce import ProgressCallback
from .ollama_client import OllamaClient
from .plans import (
    PlannedSlide,
    SlidePlan,
    build_fallback_slide_plan,
    coerce_slide_plan_payload,
    sanitize_slide_plan,
)
from .prompts import (
    CONCLUSION_FROM_KB_PROMPT,
    SLIDE_BULLETS_FROM_ATOMS_PROMPT,
    SLIDE_PLAN_PROMPT,
    SYSTEM_EXPERT_ES,
)

logger = logging.getLogger(__name__)


# ----------------------------------------------------- modelos de salida --

class SlideBullets(BaseModel):
    bullets: list[str] = Field(min_length=1, max_length=MAX_BULLETS_PER_SLIDE + 2)


@dataclass
class BuiltSlide:
    title: str
    bullets: list[str]
    kind: str = "outlook"


@dataclass
class PresentationPlan:
    title: str
    outline: list[str]
    slides: list[BuiltSlide] = field(default_factory=list)
    conclusion: list[str] = field(default_factory=list)


# ---------------------------------------------------- limpieza de bullets --

# Patrones agrupados: cualquier match descarta el bullet entero.
# - Prefijos de átomo (el LLM copió el bloque de contexto).
# - Bloques tipo "· definición ·" que sólo aparecen en el prompt.
# - Flechas internas de Relation (tripleta `—[kind]→`).
# - Anglicismos frecuentes que degradan la calidad.
_BULLET_REJECT_RE = re.compile(
    r"(?i)"
    r"(?:^\s*(?:def|ex|fc|rel|dt):)"
    r"|·\s*(?:definici[oó]n|ejemplo|dato|relaci[oó]n|f[oó]rmula|code|c[oó]digo)\s*·"
    r"|(?:—|--|-|–)\s*\[[^\]]+\]\s*(?:→|->|—>|–>)"
    r"|\b(?:blueprint|inheritance|overriding|concretad[oa])\b"
)

# Aperturas meta (bullet empieza así) — didactismo vacío.
_META_OPENINGS = (
    "en esta diapositiva", "en este apartado",
    "en esta presentacion", "en esta presentación",
    "a continuacion", "a continuación",
    "se habla de", "se trata de",
    "es importante destacar",
    "visión panorámica", "vision panoramica",
)
# Fragmentos meta en cualquier posición.
_META_EMBEDDED = (
    "sirve de ejemplo para ilustrar", "sirve para ilustrar",
    "facilita la creación y gestión", "facilita la creacion y gestion",
    "facilita la comprensión de", "facilita la comprension de",
    "ayuda a comprender", "permite entender", "permite comprender",
    "es clave para entender", "es clave para comprender",
    "ilustra cómo", "ilustra como",
)


def _echoes_title(bullet_norm: str, title_norm: str) -> bool:
    """El bullet empieza repitiendo el título + separador (ruido típico)."""
    if len(title_norm) < 6 and len(title_norm.split()) < 2:
        return False
    if not bullet_norm.startswith(title_norm):
        return False
    tail = bullet_norm[len(title_norm):].lstrip()
    return not tail or tail[:1] in {":", ".", "—", "–", "-", ",", "·"}


def _trim_by_words(text: str, max_len: int) -> str | None:
    """Recorta solo en cierre natural (. ; : ? !); no fabrica cierres."""
    if len(text) <= max_len:
        return text
    cut = text[:max_len]
    close = max(cut.rfind("."), cut.rfind(";"), cut.rfind(":"), cut.rfind("?"), cut.rfind("!"))
    if close < int(max_len * 0.45):
        return None
    out = cut[: close + 1].strip()
    return out or None


_BAD_END_RE = re.compile(
    r"(?i)\b(?:de|del|la|el|los|las|y|o|u|en|con|por|para|un|una|unos|unas|que|como)\.$"
)
_CONTENT_TOKEN_RE = re.compile(r"[a-z0-9áéíóúüñ]{3,}", flags=re.IGNORECASE)
_SENTENCE_END_RE = re.compile(r"[.;:?!]\)?\s*$")
_STRIP_PUNCT_RE = re.compile(r"[^\w\sáéíóúüñ]+", flags=re.IGNORECASE)
_BULLET_STOPWORDS = {
    "de", "del", "la", "el", "los", "las", "un", "una", "unos", "unas", "y", "o", "u",
    "en", "con", "por", "para", "que", "como", "se", "es", "son", "al", "lo", "su",
}


def _normalize_bullet_ending(text: str) -> str | None:
    """Valida cierre natural y descarta finales sospechosos."""
    t = text.strip()
    if not t:
        return None
    t = t.rstrip(",·-—").strip()
    if not t or not _SENTENCE_END_RE.search(t):
        return None
    if _BAD_END_RE.search(t):
        return None
    return t if len(t) >= 25 else None


def _balance_bullet_density(bullets: list[str]) -> list[str]:
    """Equilibra densidad visual: más bullets -> bullets más cortos."""
    if not bullets:
        return []

    dedup: list[str] = []
    seen: set[str] = set()
    for b in bullets:
        key = " ".join(b.lower().split())
        if key not in seen:
            seen.add(key)
            dedup.append(b)
    out = dedup[:MAX_BULLETS_PER_SLIDE]

    max_chars = 170 if len(out) <= 3 else 135 if len(out) == 4 else 115
    balanced: list[str] = []
    for b in out:
        trimmed = _trim_by_words(b, max_chars)
        if trimmed is None:
            continue
        c = _normalize_bullet_ending(trimmed)
        if c:
            balanced.append(c)
    return balanced


def _slide_text_signature(slide: BuiltSlide) -> set[str]:
    text = " ".join(slide.bullets).lower()
    return {t for t in _CONTENT_TOKEN_RE.findall(text)}


def _normalize_for_overlap(text: str) -> set[str]:
    n = unicodedata.normalize("NFKD", text.lower())
    n = "".join(c for c in n if not unicodedata.combining(c))
    n = _STRIP_PUNCT_RE.sub(" ", n)
    return {t for t in n.split() if len(t) >= 3 and t not in _BULLET_STOPWORDS}


def _is_similar_bullet(a: str, b: str, *, threshold: float = 0.72) -> bool:
    ta, tb = _normalize_for_overlap(a), _normalize_for_overlap(b)
    if not ta or not tb:
        return False
    return (len(ta & tb) / len(ta | tb)) >= threshold


def _dedupe_bullets(bullets: list[str], prior: list[str]) -> list[str]:
    out: list[str] = []
    for b in bullets:
        if any(_is_similar_bullet(b, x) for x in out):
            continue
        if any(_is_similar_bullet(b, x) for x in prior):
            continue
        out.append(b)
    return out


def _is_near_duplicate_slide(candidate: BuiltSlide, existing: list[BuiltSlide]) -> bool:
    cand_sig = _slide_text_signature(candidate)
    if not cand_sig:
        return False
    for prev in existing:
        prev_sig = _slide_text_signature(prev)
        if not prev_sig:
            continue
        jac = len(cand_sig & prev_sig) / len(cand_sig | prev_sig)
        if jac >= 0.78:
            return True
    return False


def _is_intro_title(text: str) -> bool:
    n = unicodedata.normalize("NFKD", (text or "").lower())
    n = "".join(c for c in n if not unicodedata.combining(c))
    n = " ".join(_STRIP_PUNCT_RE.sub(" ", n).split())
    return n == "introduccion" or n.startswith("introduccion ")


def _ensure_intro_first(plan: SlidePlan, kb: KnowledgeBase) -> SlidePlan:
    """Garantiza la estructura: Introducción -> desarrollo.

    Si existe una slide de introducción, la mueve al inicio. Si no existe,
    inserta una intro mínima y recorta al máximo permitido.
    """
    slides = list(plan.slides)
    intro_idx = next(
        (i for i, s in enumerate(slides) if s.kind == "intro" or _is_intro_title(s.title)),
        None,
    )
    if intro_idx is None:
        intro = PlannedSlide(
            title="Introducción",
            kind="intro",
            atom_ids=[],
            focus=f"Panorámica general sobre {kb.main_topic}."[:200],
        )
        slides.insert(0, intro)
    elif intro_idx != 0:
        intro = slides.pop(intro_idx)
        slides.insert(0, intro)
    if len(slides) > DEFAULT_NUM_SLIDES_MAX:
        slides = slides[:DEFAULT_NUM_SLIDES_MAX]
    return SlidePlan(presentation_title=plan.presentation_title, slides=slides)


def _truncate_title(text: str, max_len: int) -> str:
    text = " ".join(text.split())
    return text if len(text) <= max_len else text[:max_len - 1].rstrip() + "…"


def _clean_bullet(text: str, max_len: int, slide_title: str | None = None) -> str | None:
    """Normaliza y valida un bullet. Devuelve `None` si hay que descartarlo."""
    if not isinstance(text, str):
        return None
    t = " ".join(text.split()).strip()
    if not t or "…" in t or "..." in t or len(t) < 25:
        return None
    if _BULLET_REJECT_RE.search(t):
        return None
    low = t.lower()
    if any(low.startswith(m) for m in _META_OPENINGS):
        return None
    if any(p in low for p in _META_EMBEDDED):
        return None
    if slide_title:
        title_norm = " ".join(slide_title.split()).strip().lower()
        if title_norm and _echoes_title(low, title_norm):
            return None
    # En modelos pequeños, bullets muy largos se truncan peor; forzamos
    # una longitud más conservadora para evitar finales rotos.
    eff_max = min(max_len, 170)
    was_trimmed = len(t) > eff_max
    if was_trimmed and not re.search(r"[.!?]\s", t[:eff_max]):
        return None
    trimmed = _trim_by_words(t, eff_max)
    if trimmed is None:
        return None
    return _normalize_bullet_ending(trimmed)


# ------------------------------------------------- bloques de contexto ----

def _atom_block(
    atom: Definition | Example | FormulaOrCode | NumericDatum | Relation,
) -> str:
    """Renderiza un átomo como bloque Markdown compacto para prompts."""
    if isinstance(atom, Definition):
        tag = " (literal)" if atom.verbatim else ""
        return f"- `{atom.id}` · definición{tag} · **{atom.term}**: {atom.definition}"
    if isinstance(atom, Example):
        attrs = ", ".join(atom.attributes) or "—"
        methods = ", ".join(atom.methods) or "—"
        return (
            f"- `{atom.id}` · ejemplo · **{atom.name}**: {atom.description}\n"
            f"  - atributos: {attrs}\n  - métodos: {methods}"
        )
    if isinstance(atom, FormulaOrCode):
        caption = atom.caption or "—"
        lang = atom.language or ""
        return f"- `{atom.id}` · {atom.kind} · {caption}\n  ```{lang}\n  {atom.content}\n  ```"
    if isinstance(atom, NumericDatum):
        return f"- `{atom.id}` · dato · **{atom.value}**: {atom.description}"
    if isinstance(atom, Relation):
        return f"- `{atom.id}` · relación · {relation_to_natural(atom)}"
    return f"- `{getattr(atom, 'id', '?')}` · (tipo desconocido)"


def _atoms_for_slide(kb: KnowledgeBase, slide: PlannedSlide) -> str:
    if not slide.atom_ids:
        return "(sin átomos asignados; genera bullets panorámicos coherentes con el tipo)"
    parts: list[str] = []
    for aid in slide.atom_ids:
        atom = kb.get_atom(aid)
        if atom is None:
            logger.warning("Átomo '%s' no existe en KB durante render.", aid)
            continue
        parts.append(_atom_block(atom))
    return "\n".join(parts) or "(sin átomos válidos)"


# --------------------------------------------------- planificación slides --

def plan_slides(client: OllamaClient, kb: KnowledgeBase) -> SlidePlan:
    """Pide al LLM el `SlidePlan` con 3 capas de resiliencia.

    Coerción defensiva → reintento con temperatura baja → fallback
    determinístico desde la KB (para modelos pequeños poco cooperativos).
    """
    if kb.atom_count == 0:
        raise GenerationError("La KB no contiene átomos; no se puede planificar la presentación.")

    prompt = SLIDE_PLAN_PROMPT.format(kb_context=kb.to_prompt_context(max_chars=6000))

    def _try(raw: Any) -> SlidePlan | None:
        data = coerce_slide_plan_payload(raw, kb)
        if data is None:
            return None
        try:
            return SlidePlan(**data)
        except ValidationError as exc:
            logger.warning("SlidePlan ValidationError: %s", exc.errors()[:2])
            return None

    plan = _try(client.generate_json(prompt, system=SYSTEM_EXPERT_ES, temperature=0.2))
    if plan is None:
        logger.warning("SlidePlan inicial inválido; reintentando con temperatura 0.1…")
        plan = _try(client.generate_json(prompt, system=SYSTEM_EXPERT_ES, temperature=0.1))
    if plan is None:
        logger.warning("SlidePlan tampoco válido; usando fallback determinístico desde la KB.")
        plan = build_fallback_slide_plan(kb)

    plan = sanitize_slide_plan(plan, kb)

    # Complemento si el saneado dejó el plan corto.
    if len(plan.slides) < DEFAULT_NUM_SLIDES_MIN:
        logger.warning("SlidePlan corto (%d<%d); complementando con fallback.",
                       len(plan.slides), DEFAULT_NUM_SLIDES_MIN)
        fb = build_fallback_slide_plan(kb)
        titles = {s.title.lower() for s in plan.slides}
        merged = list(plan.slides)
        for s in fb.slides:
            if s.title.lower() not in titles and len(merged) < 20:
                merged.append(s)
                titles.add(s.title.lower())
        plan = sanitize_slide_plan(
            SlidePlan(presentation_title=plan.presentation_title, slides=merged), kb
        )

    if len(plan.slides) > DEFAULT_NUM_SLIDES_MAX:
        plan = SlidePlan(presentation_title=plan.presentation_title,
                         slides=plan.slides[:DEFAULT_NUM_SLIDES_MAX])

    plan = _ensure_intro_first(plan, kb)

    plan.presentation_title = _truncate_title(plan.presentation_title, 120)
    for s in plan.slides:
        s.title = _truncate_title(s.title, MAX_CHARS_SLIDE_TITLE)
    logger.info("SlidePlan: '%s', %d slides", plan.presentation_title, len(plan.slides))
    return plan


# ------------------------------------------- generación de bullets (LLM) --

def _render_bullets(
    client: OllamaClient,
    prompt: str,
    *,
    slide_title: str,
    temperature: float = 0.3,
    empty_raises: bool = True,
) -> list[str]:
    """Renderiza bullets: llama al LLM, valida JSON y limpia.

    Compartido por slides de contenido y por la conclusión. Si no sobrevive
    ningún bullet a la limpieza y `empty_raises=True`, lanza `GenerationError`
    para que el caller decida (regenerar o marcar fallo).
    """
    raw = client.generate_json(prompt, system=SYSTEM_EXPERT_ES, temperature=temperature)
    if not isinstance(raw, dict):
        raise GenerationError(f"Bullets inválidos para '{slide_title}' (no es JSON objeto).")
    try:
        sb = SlideBullets(**raw)
    except ValidationError as exc:
        raise GenerationError(
            f"Bullets con estructura inválida ({slide_title}): {exc.errors()[:1]}"
        ) from exc

    out = [c for b in sb.bullets
           if (c := _clean_bullet(b, MAX_CHARS_PER_BULLET, slide_title=slide_title))]
    out = _balance_bullet_density(out)
    if not out and empty_raises:
        raise GenerationError(
            f"Todos los bullets descartados por calidad en slide '{slide_title}'."
        )
    return out[:MAX_BULLETS_PER_SLIDE]


def render_slide_bullets(
    client: OllamaClient,
    kb: KnowledgeBase,
    slide: PlannedSlide,
    plan: SlidePlan,
    index: int,
    total: int,
    prior_bullets: list[str],
) -> list[str]:
    prompt = SLIDE_BULLETS_FROM_ATOMS_PROMPT.format(
        presentation_title=plan.presentation_title,
        index=index, total=total,
        slide_title=slide.title,
        kind=slide.kind,
        focus=(slide.focus or "—"),
        atom_details=_atoms_for_slide(kb, slide),
        outline="\n".join(f"- {s.title}" for s in plan.slides),
        anti_repeat_context=(
            "\n".join(f"- {b}" for b in prior_bullets[-8:]) if prior_bullets else "(sin contenido previo)"
        ),
    )
    first = _dedupe_bullets(
        _render_bullets(client, prompt, slide_title=slide.title, temperature=0.3),
        prior_bullets,
    )
    if len(first) >= 3:
        return first

    retry_prompt = (
        prompt
        + "\n\nREINTENTO ÚNICO:\n"
          "- Debes devolver 3-4 bullets completos y cerrados.\n"
          "- No repitas ideas ya cubiertas en el contenido previo.\n"
          "- Si no cabe una frase completa, omítela.\n"
    )
    retry = _dedupe_bullets(
        _render_bullets(
            client, retry_prompt, slide_title=slide.title, temperature=0.2, empty_raises=False
        ),
        prior_bullets,
    )
    return retry if retry else first


def _render_conclusion(
    client: OllamaClient, kb: KnowledgeBase, presentation_title: str
) -> list[str]:
    prompt = CONCLUSION_FROM_KB_PROMPT.format(
        presentation_title=presentation_title,
        kb_context=kb.to_prompt_context(max_chars=4000),
    )
    return _render_bullets(
        client, prompt, slide_title="Conclusiones",
        temperature=0.3, empty_raises=False,
    )


# ----------------------------------------- construcción determinista PPTX --

def _set_title(shape: Any, text: str) -> None:
    if shape is None or not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.clear()
    tf.text = text


def _replace_text_preserving_format(shape: Any, new_text: str) -> bool:
    """Sustituye el texto preservando el formato del primer run.

    Necesario en shapes que NO son placeholders (texto libre de plantillas
    importadas): `tf.clear()` pierde fuente, tamaño y color. Aquí reusamos
    el primer run como "maestro" y vaciamos el resto.
    """
    if shape is None or not shape.has_text_frame:
        return False
    tf = shape.text_frame
    if not tf.paragraphs:
        tf.text = new_text
        return True

    runs = list(tf.paragraphs[0].runs)
    if runs:
        runs[0].text = new_text
        for extra in runs[1:]:
            extra.text = ""
    else:
        tf.paragraphs[0].text = new_text

    for extra_p in list(tf.paragraphs[1:]):
        for r in list(extra_p.runs):
            r.text = ""
    return True


def _shape_font_size_pt(shape: Any) -> float:
    if not shape.has_text_frame:
        return 0.0
    return max(
        (r.font.size.pt for p in shape.text_frame.paragraphs for r in p.runs
         if r.font.size is not None),
        default=0.0,
    )


def _overwrite_template_cover(slide: Any, title: str, subtitle: str | None = None) -> None:
    """Sustituye portada de plantilla conservando su formato.

    Identificación del shape de título: (1) por marcador textual
    ("TÍTULO"/"SUBTÍTULO") o (2) por tamaño de fuente más grande.
    """
    shapes = [sh for sh in slide.shapes
              if getattr(sh, "has_text_frame", False) and not sh.is_placeholder]
    if not shapes:
        logger.warning("Portada de plantilla sin shapes de texto editables.")
        return

    title_markers = {"titulo", "título", "title"}
    subtitle_markers = {"subtitulo", "subtítulo", "subtitle"}
    title_shape = subtitle_shape = None

    for sh in shapes:
        text = " ".join(sh.text_frame.text.split()).strip().lower()
        if not title_shape and text in title_markers:
            title_shape = sh
        elif not subtitle_shape and text in subtitle_markers:
            subtitle_shape = sh

    if title_shape is None or (subtitle and subtitle_shape is None):
        by_size = sorted(shapes, key=_shape_font_size_pt, reverse=True)
        if title_shape is None and by_size:
            title_shape = by_size[0]
        if subtitle and subtitle_shape is None:
            subtitle_shape = next((sh for sh in by_size if sh is not title_shape), None)

    if title_shape is not None:
        _replace_text_preserving_format(title_shape, title)
    else:
        logger.warning("No se pudo localizar el shape del título en la portada.")

    if subtitle and subtitle_shape is not None:
        _replace_text_preserving_format(subtitle_shape, subtitle)


def _overwrite_template_body_slide(slide: Any, title: str, bullets: list[str]) -> bool:
    """Rellena una slide precargada con layout TITLE_AND_BODY."""
    title_ph = slide.shapes.title
    body_ph = _find_placeholder(slide, (1, 2, 13, 14))
    if title_ph is None or body_ph is None:
        return False
    _set_title(title_ph, title)
    _set_bullets(body_ph, bullets)
    return True


def _set_bullets(shape: Any, bullets: list[str]) -> None:
    if shape is None or not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    try:
        tf.auto_size = 1  # MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:  # noqa: BLE001
        pass
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        for run in p.runs:
            if run.font.size is None:
                run.font.size = Pt(20)


def _find_placeholder(slide: Any, idx_candidates: tuple[int, ...]) -> Any | None:
    by_idx: dict[int, Any] = {p.placeholder_format.idx: p for p in slide.placeholders}
    for idx in idx_candidates:
        if idx in by_idx:
            return by_idx[idx]
    return next((p for p in slide.placeholders if p.placeholder_format.idx != 0), None)


def _add_title_slide(prs: Any, title: str, subtitle: str | None = None) -> None:
    if LAYOUT_TITLE >= len(prs.slide_layouts):
        raise TemplateError(f"La plantilla no tiene el layout {LAYOUT_TITLE} (Portada).")
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
    _set_title(slide.shapes.title, title)
    if subtitle:
        ph = _find_placeholder(slide, (1, 2, 10, 11))
        if ph is not None:
            _set_bullets(ph, [subtitle])


def _add_content_slide(prs: Any, title: str, bullets: list[str]) -> None:
    if LAYOUT_CONTENT >= len(prs.slide_layouts):
        raise TemplateError(f"La plantilla no tiene el layout {LAYOUT_CONTENT} (Contenido).")
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])
    _set_title(slide.shapes.title, title)
    body = _find_placeholder(slide, (1, 2, 13, 14))
    if body is None:
        raise TemplateError("El layout de contenido no tiene un placeholder para el cuerpo.")
    _set_bullets(body, bullets)


def _load_template() -> Any:
    if not TEMPLATE_PATH.exists():
        raise TemplateError(
            f"No se encuentra la plantilla en {TEMPLATE_PATH}. "
            "Debe existir 'plantilla_universidad.pptx' en la raíz del proyecto."
        )
    try:
        return Presentation(str(TEMPLATE_PATH))
    except Exception as exc:
        raise TemplateError(f"No se pudo abrir la plantilla PPTX: {exc}") from exc


# --------------------------------------------------------- API pública ---

def build_plan(
    client: OllamaClient,
    kb: KnowledgeBase,
    progress_cb: ProgressCallback | None = None,
    *,
    refine: bool = True,
) -> PresentationPlan:
    """Planifica slides, genera bullets por slide, refina y concluye."""
    if progress_cb:
        progress_cb("slides", 0, 0, "Planificando la presentación…")
    plan = plan_slides(client, kb)

    total = len(plan.slides)
    built: list[BuiltSlide] = []
    accepted_bullets: list[str] = []
    for i, ps in enumerate(plan.slides, start=1):
        if progress_cb:
            progress_cb("slides", i, total,
                        f"Redactando slide {i}/{total} ({ps.kind}): {ps.title}")
        try:
            bullets = render_slide_bullets(
                client, kb, ps, plan, index=i, total=total, prior_bullets=accepted_bullets
            )
        except GenerationError as exc:
            logger.warning("Fallo al generar slide '%s': %s", ps.title, exc)
            bullets = []
        candidate = BuiltSlide(title=ps.title, bullets=bullets, kind=ps.kind)
        if bullets and _is_near_duplicate_slide(candidate, built):
            logger.info("Slide duplicada por contenido descartada: '%s'", ps.title)
            continue
        built.append(candidate)
        accepted_bullets.extend(bullets)

    if refine:
        if progress_cb:
            progress_cb("slides", total, total, "Revisando calidad y puliendo slides…")
        from .critics import refine_slides as _refine_slides  # import diferido
        built, review = _refine_slides(client, kb, built, plan)
        if review.issues:
            logger.info("Revisor slides: %d issues tras refinamiento (bloqueantes: %d).",
                        len(review.issues), len(review.blocker_indices()))

    if progress_cb:
        progress_cb("slides", total, total, "Generando conclusiones…")
    conclusion = _render_conclusion(client, kb, plan.presentation_title)

    return PresentationPlan(
        title=plan.presentation_title,
        outline=[s.title for s in plan.slides],
        slides=built,
        conclusion=conclusion,
    )


def render_pptx(plan: PresentationPlan, output_path: Path | None = None) -> bytes:
    """Monta el PPTX reusando slides precargadas de la plantilla."""
    prs = _load_template()

    slides_ok = [s for s in plan.slides if s.bullets]
    for omitted in (s for s in plan.slides if not s.bullets):
        logger.warning("Slide '%s' sin bullets; se omite en el PPTX.", omitted.title)

    index_bullets = [f"{i}. {s.title}" for i, s in enumerate(slides_ok, start=1)
                     ][: MAX_BULLETS_PER_SLIDE * 2]
    preloaded = list(prs.slides)

    if preloaded:
        _overwrite_template_cover(preloaded[0], plan.title,
                                  subtitle="Presentación generada automáticamente")
    else:
        _add_title_slide(prs, plan.title, subtitle="Presentación generada automáticamente")

    index_written = False
    if len(preloaded) >= 2:
        index_written = _overwrite_template_body_slide(preloaded[1], "Índice", index_bullets)
        if not index_written:
            logger.warning("Slide 2 de plantilla sin TITLE+BODY; índice como slide nueva.")
    if not index_written:
        _add_content_slide(prs, "Índice", index_bullets)

    for slide in slides_ok:
        _add_content_slide(prs, slide.title, slide.bullets)
    if plan.conclusion:
        _add_content_slide(prs, "Conclusiones", plan.conclusion)

    buffer = BytesIO()
    prs.save(buffer)
    data = buffer.getvalue()
    if output_path is not None:
        output_path.write_bytes(data)
    return data


def generate_presentation(
    client: OllamaClient,
    kb: KnowledgeBase,
    progress_cb: ProgressCallback | None = None,
    output_path: Path | None = None,
) -> tuple[bytes, PresentationPlan]:
    """Pipeline completo: KB -> plan -> bullets -> PPTX bytes."""
    plan = build_plan(client, kb, progress_cb=progress_cb)
    data = render_pptx(plan, output_path=output_path)
    return data, plan
